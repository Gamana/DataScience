"""
Export Retail Sales Dashboard to PowerPoint (.pptx)
Matches the 3-page interactive dashboard: Executive, Product, Geographic.
"""

from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


ROOT = Path(__file__).parent
DATA = ROOT / "data"
OUT_DIR = ROOT / "export"
CHART_DIR = OUT_DIR / "chart_images"
PPTX_PATH = OUT_DIR / "Retail_Sales_Dashboard.pptx"

# Corporate palette
BLUE = "0078D4"
DARK = "252525"
GRAY = "505050"
LIGHT = "F5F7FA"
WHITE = "FFFFFF"
GOOD = "107C10"
TEAL = "00B7C3"
BORDER = "E1E1E1"

BLUE_RGB = rgb(BLUE)
DARK_RGB = rgb(DARK)
GRAY_RGB = rgb(GRAY)
WHITE_RGB = rgb(WHITE)
GOOD_RGB = rgb(GOOD)


def money(n: float) -> str:
    if abs(n) >= 1_000_000:
        return f"${n/1_000_000:.2f}M"
    if abs(n) >= 1_000:
        return f"${n/1_000:.1f}K"
    return f"${n:,.0f}"


def set_run_font(run, size=12, bold=False, color=DARK_RGB, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=DARK_RGB, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill_hex, line_hex=None, rounded=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
    else:
        shape.line.fill.background()
    return shape


def style_chart(chart, show_legend=True):
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = False


def load_data():
    flat = pd.read_csv(DATA / "Retail_Sales_Flat.csv", parse_dates=["Order Date"])
    geo = pd.read_csv(DATA / "Dim_Geography.csv")
    pipeline = pd.read_csv(DATA / "Dim_Sales_Pipeline.csv")
    targets = pd.read_csv(DATA / "Dim_Targets.csv")
    flat = flat.merge(
        geo[["City", "State", "Latitude", "Longitude", "Region"]],
        on=["City", "State", "Region"],
        how="left",
    )
    flat["Year"] = flat["Order Date"].dt.year
    flat["YearMonth"] = flat["Order Date"].dt.strftime("%Y-%m")
    return flat, pipeline, targets


def make_matplotlib_charts(flat, pipeline):
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "axes.edgecolor": "#E1E1E1",
        "axes.labelcolor": "#505050",
        "xtick.color": "#505050",
        "ytick.color": "#505050",
        "figure.facecolor": "white",
        "axes.facecolor": "white",
    })

    paths = {}

    # --- Scatter ---
    prod = flat.groupby("Product Name", as_index=False).agg(
        Sales=("Sales", "sum"), Profit=("Profit", "sum"), Quantity=("Quantity", "sum")
    )
    fig, ax = plt.subplots(figsize=(7.2, 4.2))
    sizes = np.sqrt(prod["Quantity"]) * 8
    ax.scatter(prod["Sales"], prod["Profit"], s=sizes, c="#0078D4", alpha=0.55, edgecolors="#005A9E", linewidths=0.5)
    ax.set_xlabel("Sales ($)")
    ax.set_ylabel("Profit ($)")
    ax.set_title("Sales vs Profit (bubble size = Quantity)", fontsize=11, color="#252525", pad=10)
    ax.axhline(0, color="#D83B01", linewidth=0.8, linestyle="--", alpha=0.5)
    ax.grid(True, alpha=0.25)
    fig.tight_layout()
    paths["scatter"] = CHART_DIR / "scatter.png"
    fig.savefig(paths["scatter"], dpi=160, bbox_inches="tight")
    plt.close(fig)

    # --- Treemap-like ---
    by_cat = flat.groupby("Product Category")["Sales"].sum().sort_values(ascending=False)
    colors = {"Technology": "#0078D4", "Furniture": "#505050", "Office Supplies": "#00B7C3"}
    fig, ax = plt.subplots(figsize=(7.2, 3.6))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    ax.set_title("Sales by Category (Tree Map)", fontsize=11, color="#252525", pad=8)
    total = by_cat.sum()
    x = 0
    for cat, val in by_cat.items():
        w = val / total
        rect = mpatches.FancyBboxPatch(
            (x, 0.15), w - 0.008, 0.7,
            boxstyle="round,pad=0.01,rounding_size=0.02",
            facecolor=colors.get(cat, "#8764B8"),
            edgecolor="white",
            linewidth=2,
        )
        ax.add_patch(rect)
        ax.text(x + w / 2, 0.55, cat, ha="center", va="center", color="white", fontsize=10, fontweight="bold")
        ax.text(x + w / 2, 0.38, money(val), ha="center", va="center", color="white", fontsize=9)
        x += w
    fig.tight_layout()
    paths["treemap"] = CHART_DIR / "treemap.png"
    fig.savefig(paths["treemap"], dpi=160, bbox_inches="tight")
    plt.close(fig)

    # --- Funnel ---
    pipe = pipeline.sort_values("Stage Order")
    fig, ax = plt.subplots(figsize=(5.5, 4.2))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, len(pipe) + 0.5)
    ax.axis("off")
    ax.set_title("Sales Pipeline Funnel", fontsize=11, color="#252525", pad=8)
    shades = ["#0078D4", "#2B88D8", "#00B7C3", "#505050", "#107C10"]
    max_v = pipe["Pipeline Value"].max()
    for i, (_, r) in enumerate(pipe.iterrows()):
        width = 0.35 + 0.55 * (r["Pipeline Value"] / max_v)
        y = len(pipe) - i - 0.8
        left = (1 - width) / 2
        rect = mpatches.FancyBboxPatch(
            (left, y), width, 0.7,
            boxstyle="round,pad=0.01,rounding_size=0.02",
            facecolor=shades[i % len(shades)],
            edgecolor="white",
        )
        ax.add_patch(rect)
        ax.text(0.5, y + 0.35, f"{r['Stage']}  ·  {money(r['Pipeline Value'])}",
                ha="center", va="center", color="white", fontsize=9, fontweight="bold")
    fig.tight_layout()
    paths["funnel"] = CHART_DIR / "funnel.png"
    fig.savefig(paths["funnel"], dpi=160, bbox_inches="tight")
    plt.close(fig)

    # --- City bubble map (simplified US scatter) ---
    city = flat.groupby(["City", "Latitude", "Longitude"], as_index=False)["Sales"].sum()
    fig, ax = plt.subplots(figsize=(7.5, 4.5))
    sizes = (city["Sales"] / city["Sales"].max()) * 400 + 20
    ax.scatter(city["Longitude"], city["Latitude"], s=sizes, c="#0078D4", alpha=0.45, edgecolors="#005A9E")
    ax.set_xlabel("Longitude")
    ax.set_ylabel("Latitude")
    ax.set_title("Map — Sales by City", fontsize=11, color="#252525", pad=10)
    ax.set_xlim(-130, -65)
    ax.set_ylim(24, 50)
    ax.grid(True, alpha=0.2)
    # label top 5 cities
    top = city.nlargest(5, "Sales")
    for _, r in top.iterrows():
        ax.annotate(r["City"], (r["Longitude"], r["Latitude"]), fontsize=7, color="#252525",
                    xytext=(4, 4), textcoords="offset points")
    fig.tight_layout()
    paths["map"] = CHART_DIR / "map_city.png"
    fig.savefig(paths["map"], dpi=160, bbox_inches="tight")
    plt.close(fig)

    # --- Filled state choropleth (bar-style state ranking as visual stand-in + colored bars) ---
    state = flat.groupby("State")["Sales"].sum().sort_values(ascending=True)
    fig, ax = plt.subplots(figsize=(7.5, 5.2))
    norm = state / state.max()
    cmap_colors = [plt.cm.Blues(0.35 + 0.65 * v) for v in norm]
    ax.barh(state.index, state.values, color=cmap_colors)
    ax.set_xlabel("Sales ($)")
    ax.set_title("Filled Map Proxy — Sales by State", fontsize=11, color="#252525", pad=10)
    ax.tick_params(axis="y", labelsize=7)
    fig.tight_layout()
    paths["filled"] = CHART_DIR / "filled_state.png"
    fig.savefig(paths["filled"], dpi=160, bbox_inches="tight")
    plt.close(fig)

    # --- Area profit ---
    monthly = flat.groupby("YearMonth", as_index=False).agg(Profit=("Profit", "sum"), Quantity=("Quantity", "sum"), Sales=("Sales", "sum"))
    monthly = monthly.sort_values("YearMonth")
    # sample every month but label sparsely
    fig, ax = plt.subplots(figsize=(8.5, 3.8))
    ax.fill_between(range(len(monthly)), monthly["Profit"], color="#107C10", alpha=0.25)
    ax.plot(range(len(monthly)), monthly["Profit"], color="#107C10", linewidth=2)
    ax.set_title("Monthly Profit (Area)", fontsize=11, color="#252525")
    ax.set_ylabel("Profit ($)")
    ticks = list(range(0, len(monthly), max(1, len(monthly) // 8)))
    ax.set_xticks(ticks)
    ax.set_xticklabels([monthly.iloc[i]["YearMonth"] for i in ticks], rotation=30, ha="right", fontsize=8)
    ax.grid(True, alpha=0.25)
    fig.tight_layout()
    paths["area"] = CHART_DIR / "area_profit.png"
    fig.savefig(paths["area"], dpi=160, bbox_inches="tight")
    plt.close(fig)

    return paths, monthly, prod, by_cat


def add_kpi_card(slide, left, top, width, height, title, value, subtitle=None):
    add_rect(slide, left, top, width, height, WHITE, BORDER)
    add_textbox(slide, left + 0.12, top + 0.08, width - 0.2, 0.28, title, size=10, bold=True, color=GRAY_RGB)
    add_textbox(slide, left + 0.12, top + 0.32, width - 0.2, 0.45, value, size=22, bold=True, color=BLUE_RGB)
    if subtitle:
        add_textbox(slide, left + 0.12, top + 0.75, width - 0.2, 0.28, subtitle, size=9, color=GRAY_RGB)


def build_pptx(flat, pipeline, targets, paths, monthly, prod, by_cat):
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    total_sales = flat["Sales"].sum()
    total_profit = flat["Profit"].sum()
    total_orders = flat["Order ID"].nunique()
    total_customers = flat["Customer ID"].nunique()
    margin = total_profit / total_sales
    aov = total_sales / total_orders

    sales_target = targets["Sales Target"].sum()
    profit_target = targets["Profit Target"].sum()
    sales_ach = total_sales / sales_target
    profit_ach = total_profit / profit_target

    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 1.4, BLUE)
    add_textbox(slide, 0.6, 0.35, 12, 0.6, "Retail Sales Dashboard", size=32, bold=True, color=WHITE_RGB)
    add_textbox(slide, 0.6, 0.9, 12, 0.35, "Executive · Product · Geographic Analytics", size=14, color=WHITE_RGB)
    add_textbox(
        slide, 0.6, 2.0, 11, 1.2,
        f"5,200 orders  |  {flat['Order Date'].min().date()} → {flat['Order Date'].max().date()}\n"
        f"Total Sales {money(total_sales)}  ·  Total Profit {money(total_profit)}  ·  "
        f"Profit Margin {margin*100:.1f}%",
        size=16, color=DARK_RGB,
    )
    add_textbox(
        slide, 0.6, 3.4, 11, 1.5,
        "This deck exports the interactive dashboard visuals for presentations and teaching.\n"
        "Theme: Corporate Blue / Dark Gray / White  ·  Source: Retail_Sales_StarSchema",
        size=13, color=GRAY_RGB,
    )
    add_textbox(slide, 0.6, 6.8, 12, 0.3, "Power BI Retail Dashboard Package", size=10, color=GRAY_RGB)

    # ========== SLIDE 2: Executive KPIs ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 10, 0.4, "Page 1 — Executive Dashboard · KPI Overview", size=18, bold=True, color=WHITE_RGB)

    gap = 0.22
    card_w = 2.9
    start = 0.45
    cards = [
        ("Total Sales", money(total_sales), f"AOV {money(aov)}"),
        ("Total Profit", money(total_profit), f"Margin {margin*100:.1f}%"),
        ("Total Orders", f"{total_orders:,}", "Distinct Order IDs"),
        ("Total Customers", f"{total_customers:,}", "Unique Customer IDs"),
    ]
    for i, (t, v, s) in enumerate(cards):
        add_kpi_card(slide, start + i * (card_w + gap), 1.05, card_w, 1.15, t, v, s)

    add_kpi_card(
        slide, 0.45, 2.45, 6.0, 1.25,
        "KPI — Sales vs Target",
        money(total_sales),
        f"Target {money(sales_target)}  ·  Achievement {sales_ach*100:.1f}%",
    )
    add_kpi_card(
        slide, 6.75, 2.45, 6.0, 1.25,
        "KPI — Profit vs Target",
        money(total_profit),
        f"Target {money(profit_target)}  ·  Achievement {profit_ach*100:.1f}%",
    )

    # Monthly sales line (native)
    chart_data = CategoryChartData()
    # Use quarterly or every 3rd month for readability
    m = monthly.copy()
    step = max(1, len(m) // 12)
    m_plot = m.iloc[::step]
    chart_data.categories = list(m_plot["YearMonth"])
    chart_data.add_series("Sales", [float(x) for x in m_plot["Sales"]])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.45), Inches(4.0), Inches(6.2), Inches(3.1), chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Monthly Sales Trend"
    style_chart(chart, show_legend=False)

    # Category columns
    chart_data2 = CategoryChartData()
    cats = flat.groupby("Product Category")["Sales"].sum().sort_values(ascending=False)
    chart_data2.categories = list(cats.index)
    chart_data2.add_series("Sales", [float(x) for x in cats.values])
    chart2 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.9), Inches(4.0), Inches(5.9), Inches(3.1), chart_data2
    ).chart
    chart2.has_title = True
    chart2.chart_title.text_frame.paragraphs[0].text = "Sales by Category"
    style_chart(chart2, show_legend=False)

    # ========== SLIDE 3: Executive Segment Donut ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Page 1 — Executive Dashboard · Segment Mix", size=18, bold=True, color=WHITE_RGB)

    seg = flat.groupby("Segment")["Sales"].sum().sort_values(ascending=False)
    chart_data = CategoryChartData()
    chart_data.categories = list(seg.index)
    chart_data.add_series("Sales", [float(x) for x in seg.values])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(1.2), Inches(1.3), Inches(6.5), Inches(5.5), chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Sales by Segment (Donut)"
    style_chart(chart, show_legend=True)

    add_rect(slide, 8.2, 1.5, 4.5, 4.8, WHITE, BORDER)
    add_textbox(slide, 8.4, 1.7, 4.1, 0.4, "Segment Insights", size=16, bold=True, color=DARK_RGB)
    y = 2.3
    for name, val in seg.items():
        pct = val / total_sales * 100
        add_textbox(slide, 8.4, y, 4.1, 0.55, f"{name}\n{money(val)}  ({pct:.1f}%)", size=13, color=GRAY_RGB)
        y += 0.9
    add_textbox(slide, 8.4, 5.5, 4.1, 0.6, "Slicers in live dashboard:\nYear · Category · Region", size=11, color=GRAY_RGB)

    # ========== SLIDE 4: Product Table + Top 10 ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Page 2 — Product Analysis · Top Products", size=18, bold=True, color=WHITE_RGB)

    top10 = (
        flat.groupby(["Product Name", "Product Category"], as_index=False)
        .agg(Quantity=("Quantity", "sum"), Sales=("Sales", "sum"), Profit=("Profit", "sum"))
        .sort_values("Sales", ascending=False)
        .head(10)
    )

    rows, cols = 11, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.0), Inches(6.4), Inches(5.8)).table
    headers = ["Product Name", "Category", "Qty", "Sales", "Profit"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=9, bold=True, color=WHITE_RGB)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_RGB
    for i, (_, r) in enumerate(top10.iterrows(), start=1):
        vals = [r["Product Name"], r["Product Category"], f"{int(r['Quantity']):,}", money(r["Sales"]), money(r["Profit"])]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run_font(run, size=8, color=DARK_RGB)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(LIGHT)

    chart_data = CategoryChartData()
    chart_data.categories = [p[:22] + ("…" if len(p) > 22 else "") for p in top10["Product Name"]]
    chart_data.add_series("Sales", [float(x) for x in top10["Sales"]])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.0), Inches(1.0), Inches(5.9), Inches(5.8), chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Top 10 Products by Sales"
    style_chart(chart, show_legend=False)

    # ========== SLIDE 5: Matrix ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Page 2 — Product Analysis · Matrix (Category × Year)", size=18, bold=True, color=WHITE_RGB)

    mat = flat.groupby(["Product Category", "Year"])[["Sales", "Profit"]].sum().reset_index()
    years = sorted(flat["Year"].unique())
    cats = sorted(flat["Product Category"].unique())
    # header: Category | 2022 Sales | 2022 Profit | ...
    col_headers = ["Category"]
    for y in years:
        col_headers += [f"{y} Sales", f"{y} Profit"]
    n_cols = len(col_headers)
    n_rows = len(cats) + 1
    table = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.5)).table
    for j, h in enumerate(col_headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=9, bold=True, color=WHITE_RGB)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_RGB
    for i, cat in enumerate(cats, start=1):
        table.cell(i, 0).text = cat
        for p in table.cell(i, 0).text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=10, bold=True, color=DARK_RGB)
        for k, y in enumerate(years):
            row = mat[(mat["Product Category"] == cat) & (mat["Year"] == y)]
            s = float(row["Sales"].iloc[0]) if len(row) else 0
            pft = float(row["Profit"].iloc[0]) if len(row) else 0
            table.cell(i, 1 + k * 2).text = money(s)
            table.cell(i, 2 + k * 2).text = money(pft)
            for idx in (1 + k * 2, 2 + k * 2):
                for p in table.cell(i, idx).text_frame.paragraphs:
                    for r in p.runs:
                        set_run_font(r, size=10, color=DARK_RGB)
    add_textbox(
        slide, 0.4, 6.2, 12, 0.6,
        "In Power BI: Matrix rows = Category → Sub Category, columns = Year, values = Sales & Profit (with drill-down).",
        size=11, color=GRAY_RGB,
    )

    # ========== SLIDE 6: Treemap + Scatter + Pie + Funnel ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Page 2 — Product Analysis · Advanced Visuals", size=18, bold=True, color=WHITE_RGB)

    slide.shapes.add_picture(str(paths["treemap"]), Inches(0.35), Inches(0.95), Inches(6.3), Inches(2.9))
    slide.shapes.add_picture(str(paths["scatter"]), Inches(6.8), Inches(0.95), Inches(6.1), Inches(2.9))
    slide.shapes.add_picture(str(paths["funnel"]), Inches(0.35), Inches(4.0), Inches(5.2), Inches(3.1))

    # Pie payment native
    pay = flat.groupby("Payment Mode")["Sales"].sum().sort_values(ascending=False)
    chart_data = CategoryChartData()
    chart_data.categories = list(pay.index)
    chart_data.add_series("Sales", [float(x) for x in pay.values])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(5.8), Inches(4.0), Inches(7.0), Inches(3.2), chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Sales by Payment Mode"
    style_chart(chart, show_legend=True)

    # ========== SLIDE 7: Geographic ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Page 3 — Geographic Dashboard · Maps", size=18, bold=True, color=WHITE_RGB)
    slide.shapes.add_picture(str(paths["map"]), Inches(0.3), Inches(1.0), Inches(6.4), Inches(5.8))
    slide.shapes.add_picture(str(paths["filled"]), Inches(6.8), Inches(1.0), Inches(6.2), Inches(5.8))

    # ========== SLIDE 8: Geo charts ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Page 3 — Geographic Dashboard · Trends", size=18, bold=True, color=WHITE_RGB)

    slide.shapes.add_picture(str(paths["area"]), Inches(0.35), Inches(1.0), Inches(7.6), Inches(3.0))

    reg = flat.groupby("Region")["Sales"].sum().sort_values(ascending=False)
    chart_data = CategoryChartData()
    chart_data.categories = list(reg.index)
    chart_data.add_series("Sales", [float(x) for x in reg.values])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.1), Inches(1.0), Inches(4.8), Inches(3.0), chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Sales by Region"
    style_chart(chart, show_legend=False)

    # Quantity line
    chart_data = CategoryChartData()
    m_plot = monthly.iloc[:: max(1, len(monthly) // 12)]
    chart_data.categories = list(m_plot["YearMonth"])
    chart_data.add_series("Quantity", [int(x) for x in m_plot["Quantity"]])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(0.35), Inches(4.2), Inches(12.6), Inches(2.9), chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "Monthly Quantity Sold"
    style_chart(chart, show_legend=False)

    # ========== SLIDE 9: Visual inventory ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 0, 13.333, 0.7, BLUE)
    add_textbox(slide, 0.4, 0.18, 12, 0.4, "Visual Coverage Checklist", size=18, bold=True, color=WHITE_RGB)

    visuals = [
        ("Card", "Executive — Sales, Profit, Orders, Customers"),
        ("KPI", "Executive — Sales/Profit vs Target"),
        ("Slicer", "Live HTML/Power BI — Year, Category, Region"),
        ("Line Chart", "Monthly Sales Trend · Monthly Quantity"),
        ("Column Chart", "Sales by Category · Sales by Region"),
        ("Donut Chart", "Sales by Segment"),
        ("Table", "Product metrics (Name, Category, Qty, Sales, Profit)"),
        ("Matrix", "Category × Year Sales & Profit"),
        ("Bar Chart", "Top 10 Products by Sales"),
        ("Tree Map", "Sales by Category"),
        ("Scatter Chart", "Sales vs Profit (size = Quantity)"),
        ("Pie Chart", "Sales by Payment Mode"),
        ("Funnel Chart", "Lead → Qualified → Proposal → Negotiation → Won"),
        ("Map", "Sales by City"),
        ("Filled Map", "Sales by State"),
        ("Area Chart", "Monthly Profit"),
    ]
    add_rect(slide, 0.4, 1.0, 12.5, 5.9, WHITE, BORDER)
    y = 1.15
    for i, (name, desc) in enumerate(visuals):
        col = 0 if i < 8 else 1
        row = i if i < 8 else i - 8
        x = 0.6 + col * 6.3
        yy = 1.15 + row * 0.7
        add_textbox(slide, x, yy, 6.0, 0.55, f"✓  {name}\n     {desc}", size=11, color=DARK_RGB)

    # ========== SLIDE 10: Closing ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, LIGHT)
    add_rect(slide, 0, 2.4, 13.333, 2.4, BLUE)
    add_textbox(slide, 0.6, 2.7, 12, 0.6, "Thank You", size=32, bold=True, color=WHITE_RGB, align=PP_ALIGN.CENTER)
    add_textbox(
        slide, 0.6, 3.4, 12, 0.8,
        "Interactive dashboard: PowerBI_Retail_Dashboard/dashboard/dashboard.html\n"
        "Power BI build guide: docs/02_Build_Guide.md",
        size=14, color=WHITE_RGB, align=PP_ALIGN.CENTER,
    )

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main():
    print("Loading data…")
    flat, pipeline, targets = load_data()
    print("Rendering chart images…")
    paths, monthly, prod, by_cat = make_matplotlib_charts(flat, pipeline)
    print("Building PowerPoint…")
    out = build_pptx(flat, pipeline, targets, paths, monthly, prod, by_cat)
    print(f"Saved: {out}")
    print(f"Size: {out.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
